import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

# Path definitions
TEMPLATE_PATH = r"D:\Projects\Projects\AlignMate_nba\Presentation_PPt.pptx"
OUTPUT_PATH = r"D:\Projects\Projects\AlignMate_nba\Presentation_PPt.pptx"
LOGO_TEMP_PATH = "temp_logo.png"

# Color constants
ORANGE_ACCENT = RGBColor(228, 131, 18)  # #E48312 theme color
DARK_TEXT = RGBColor(28, 28, 28)       # Charcoal text color
MUTED_TEXT = RGBColor(80, 80, 80)      # Gray subtitle color

def get_content_placeholder(slide):
    # Try finding the Content Placeholder shape (index 1)
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            return shape
    # Fallback to any placeholder with text frame
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            return shape
    return None

def populate_slide_content(slide, title, bullets):
    placeholder = get_content_placeholder(slide)
    if not placeholder:
        print(f"[ERROR] Content placeholder not found on slide for title: {title}")
        return
        
    tf = placeholder.text_frame
    tf.word_wrap = True
    tf.clear()
    
    # 1. Title Paragraph (Index 0 in TextFrame)
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Times New Roman'
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = ORANGE_ACCENT
    p_title.space_after = Pt(16)
    
    # Disable bullet on the Title paragraph using direct XML manipulation
    p_title_pr = p_title._p.get_or_add_pPr()
    p_title_pr.append(parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
    
    # 2. Body Bullets
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = 'Times New Roman'
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

def main():
    if not os.path.exists(TEMPLATE_PATH):
        print(f"[ERROR] Template PPTX file not found at {TEMPLATE_PATH}")
        return
        
    print(f"Loading presentation template from {TEMPLATE_PATH}...")
    prs = Presentation(TEMPLATE_PATH)
    
    # 1. Extract Logo Image (Picture 6) from Slide 1
    logo_extracted = False
    for shape in prs.slides[0].shapes:
        if shape.name == "Picture 6":
            with open(LOGO_TEMP_PATH, "wb") as f:
                f.write(shape.image.blob)
            logo_extracted = True
            print(f"Extracted brand logo image to {LOGO_TEMP_PATH}.")
            break
            
    if not logo_extracted:
        print("[WARNING] Could not find 'Picture 6' to extract logo. New slides will not have logo overlays.")

    # =========================================================================
    # SLIDE 1: Title & Project Abstract
    # =========================================================================
    print("Updating Slide 1: Title...")
    slide1 = prs.slides[0]
    tf1 = get_content_placeholder(slide1).text_frame
    tf1.clear()
    
    # Project Title
    p1 = tf1.paragraphs[0]
    p1.text = "Project Title: AlignMate"
    p1.font.name = 'Times New Roman'
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = ORANGE_ACCENT
    p1.space_after = Pt(12)
    p1_pr = p1._p.get_or_add_pPr()
    p1_pr.append(parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
    
    # Subtitle
    p2 = tf1.add_paragraph()
    p2.text = "AI-Powered Real-Time Posture Coach & Fitness Companion"
    p2.font.name = 'Times New Roman'
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED_TEXT
    p2.space_after = Pt(40)
    p2_pr = p2._p.get_or_add_pPr()
    p2_pr.append(parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
    
    # Blank Line
    p3 = tf1.add_paragraph()
    p3.text = ""
    p3.space_after = Pt(20)
    p3_pr = p3._p.get_or_add_pPr()
    p3_pr.append(parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
    
    # Author Info
    p4 = tf1.add_paragraph()
    p4.text = "Submitted By: Arjun Bhati & Team"
    p4.font.name = 'Times New Roman'
    p4.font.size = Pt(18)
    p4.font.bold = True
    p4.font.color.rgb = DARK_TEXT
    p4_pr = p4._p.get_or_add_pPr()
    p4_pr.append(parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))

    # =========================================================================
    # SLIDE 2: Presentation Outline
    # =========================================================================
    print("Updating Slide 2: Outline...")
    slide2 = prs.slides[1]
    outline_bullets = [
        "Slide 1: Project Title & Abstract",
        "Slide 2: Presentation Outline",
        "Slide 3: Problem Statement & Engineering Context",
        "Slide 4: Pipeline System Architecture",
        "Slide 5: Data Engineering Pipeline (Angle Calculations)",
        "Slide 6: Structured Database RAG (Retrieval-Augmented Generation)",
        "Slide 7: Containerization & Isolation (Docker)",
        "Slide 8: Automated CI/CD Pipelines (GitHub Actions)",
        "Slide 9: Monitoring HUD & Client Interface",
        "Slide 10: Technical Debugging & Future Roadmap"
    ]
    populate_slide_content(slide2, "Presentation Outline", outline_bullets)

    # =========================================================================
    # SLIDE 3: Problem Statement & Engineering Context
    # =========================================================================
    print("Updating Slide 3: Problem Statement...")
    slide3 = prs.slides[2]
    bullets3 = [
        "Sedentary Lifestyles: Proved to cause chronic cervical pain, spinal stress, and long-term posture defects.",
        "Exercise Injuries: Lack of real-time form correction leads to acute gym injuries under high loading.",
        "Engineering Constraints: Must perform skeletal tracking at >15 FPS within a standard 2GB server footprint.",
        "User Data Privacy: Zero video streaming or storage in the cloud. Coordinates are processed client-side."
    ]
    populate_slide_content(slide3, "Problem Statement & Engineering Context", bullets3)

    # =========================================================================
    # SLIDE 4: Pipeline System Architecture
    # =========================================================================
    print("Updating Slide 4: System Architecture...")
    slide4 = prs.slides[3]
    bullets4 = [
        "Client Edge Processing: User's web camera captures video. MediaPipe JS local library maps 33 joints (99 floats).",
        "Persistent Transport: Streams coordinates at 15-20 FPS using persistent, low-overhead HTML5 WebSockets.",
        "FastAPI Backend: Asynchronous WebSocket pipeline (/ws & /ws/exercise) routes coordinate frames.",
        "In-Memory Analysis: Scikit-learn model classifies posture in <1.5ms. Custom rule engines verify form metrics."
    ]
    populate_slide_content(slide4, "Pipeline System Architecture", bullets4)

    # =========================================================================
    # SLIDE 5: Data Engineering Pipeline
    # =========================================================================
    print("Updating Slide 5: Data Engineering...")
    slide5 = prs.slides[4]
    bullets5 = [
        "Landmark Vectorization: Transforms 33 body coordinate landmarks into a flat 99-feature float array.",
        "Pose Normalization: Scales coordinate arrays against shoulder-width to make tracking distance-independent.",
        "Neck Angle Geometry: Midpoint M = (Left + Right)/2. Vector V = Nose - M. Angle theta calculated via cos(theta) = -V_y / ||V||.",
        "Posture Rules: Lateral Neck Tilt flagged if theta > 15 degrees. Shoulder Imbalance flagged if left/right height diff > 0.03."
    ]
    populate_slide_content(slide5, "Data Engineering Pipeline (Angle Calculations)", bullets5)

    # =========================================================================
    # ADDING SLIDES 6 - 10 (Using layout 1 and copying Logo)
    # =========================================================================
    layout_to_use = prs.slide_layouts[1] # Title and Content
    
    additional_slides = [
        # SLIDE 6
        ("Structured Database-Driven RAG Pattern", [
            "Structured Telemetry Retrieval: Instead of raw documents, queries MySQL for profile stats (age, weight, goals) and posture history.",
            "Dynamic Prompt Augmentation: Telemetry data (average posture scores, specific faults) is injected into LangChain prompt templates.",
            "Gemini Generation: Custom augmented context is sent to Google Gemini (gemini-2.5-flash) to generate tailored workouts.",
            "Output Standardization: Gemini returns a structured JSON workout and diet plan, parsed directly by the backend for user display."
        ]),
        # SLIDE 7
        ("Containerization & Isolation (Docker)", [
            "Multi-Stage Backend Build: Dockerfile compiled using python:3.10-slim. Automatically packs OpenCV & GLib OS dependencies.",
            "Frontend Static Serving: React application compiled using Node builder and served via containerized Nginx.",
            "Database Security & Isolation: MySQL 8 service is containerized, restricting port 3306 exclusively to the internal Docker network.",
            "aaPanel Production Stack: Orchestrates containers, manages reverse proxy configurations, and handles SSL certificates."
        ]),
        # SLIDE 8
        ("Automated CI/CD Pipelines (GitHub Actions)", [
            "Automatic Triggers: Pipeline automatically runs on every push or pull request to the main branch.",
            "Backend CI Workflow: Automatically builds Python virtual environment, lint-checks code, and runs tests using Pytest.",
            "Frontend Build Verification: Installs npm packages, checks syntax, and compiles React production assets.",
            "Container Build Verification: Compiles backend and frontend Dockerfiles to ensure deployment reliability."
        ]),
        # SLIDE 9
        ("Monitoring HUD & Client Interface", [
            "Real-Time Skeleton Overlay: Renders webcam feed with green (good), yellow (drift), or red (bad) skeleton landmarks.",
            "Interactive Rep HUD: Shows live rep counts, range-of-motion angles, and posture status gauges in the browser.",
            "Speech Synthesis Coaching: Client-side non-blocking audio engine calls out corrective commands like 'Sit straight'.",
            "Analytics Dashboard: Displays 14-day history, average posture scores, and tracking duration trends."
        ]),
        # SLIDE 10
        ("Technical Debugging & Future Roadmap", [
            "Network Latency Fix: Reduced lag from 500ms to <15ms by streaming coordinate floats instead of video frames.",
            "OOM Crash Prevention: Migrated from heavy local Ollama model to Google Gemini API, saving 75% server RAM.",
            "Database Pool Scoping: Resolved SQL connection exhaustion by integrating scoped database session dependency injection.",
            "Mobile App Expansion: Future roadmap includes React Native mobile client using MoveNet on-device posture tracking."
        ])
    ]

    for idx, (title, bullets) in enumerate(additional_slides):
        slide_num = idx + 6
        print(f"Creating Slide {slide_num}: {title}...")
        new_slide = prs.slides.add_slide(layout_to_use)
        
        # Add brand logo image to new slide at exact coordinate mapping
        if logo_extracted and os.path.exists(LOGO_TEMP_PATH):
            # Coordinates from original Picture 6: Left=1580457, Top=210539, Width=9092045, Height=1158340 (in EMUs)
            new_slide.shapes.add_picture(LOGO_TEMP_PATH, 1580457, 210539, 9092045, 1158340)
            
        # Populate text contents
        populate_slide_content(new_slide, title, bullets)

    # Save presentation overwriting original template
    prs.save(OUTPUT_PATH)
    print(f"[SUCCESS] Presentation generated and saved to {OUTPUT_PATH}")
    
    # Cleanup logo temp file
    if os.path.exists(LOGO_TEMP_PATH):
        os.remove(LOGO_TEMP_PATH)

if __name__ == "__main__":
    main()
